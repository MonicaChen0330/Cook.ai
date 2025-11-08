from . import Document, Page, ExtractedImage, DocumentLoader
from pptx import Presentation
from .image_utils import image_to_base64_uri
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .ocr_utils import ocr_image_to_text

class PptxLoader(DocumentLoader):
    """A loader for Microsoft PowerPoint (.pptx) files."""

    def load(self, source: str) -> Document:
        """Reads text and extracts images from a .pptx file on a slide-by-slide basis."""
        try:
            prs = Presentation(source)
            doc_pages = []

            for i, slide in enumerate(prs.slides):
                native_text_parts = []
                page_images = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        native_text_parts.append(shape.text)
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_bytes = shape.image.blob
                            
                            ocr_text = ocr_image_to_text(image_bytes) or ""
                            base64_image_uri = image_to_base64_uri(image_bytes)

                            if base64_image_uri:
                                page_images.append(ExtractedImage(image_uri=base64_image_uri, ocr_text=ocr_text))
                        except Exception as e:
                            print(f"Warning: Could not process an image on slide {i+1}: {e}")

                native_text = "\n".join(native_text_parts)
                doc_pages.append(Page(
                    page_number=i + 1, 
                    native_text=native_text,
                    extracted_images=page_images
                ))
            
            print(f"Successfully read {len(doc_pages)} slides from {source}")
            return Document(source=source, pages=doc_pages)
        except Exception as e:
            print(f"Error reading PPTX file: {str(e)}")
            raise e
